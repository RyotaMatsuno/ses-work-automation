import subprocess
import sys

sys.stdout.reconfigure(encoding="utf-8")

IMP = r"C:\Users\ma_py\OneDrive\デスクトップ\ses_work\mail_attachment_importer"
SES = r"C:\Users\ma_py\OneDrive\デスクトップ\ses_work"

# 1. 構文チェック
r = subprocess.run(
    [sys.executable, "-m", "py_compile", "file_parser.py"],
    capture_output=True,
    encoding="utf-8",
    errors="replace",
    cwd=IMP,
)
print("py_compile:", "OK" if r.returncode == 0 else f"NG: {r.stderr}")

# 2. pptxライブラリ確認
r2 = subprocess.run(
    [sys.executable, "-c", 'from pptx import Presentation; print("pptx OK")'],
    capture_output=True,
    encoding="utf-8",
    errors="replace",
    cwd=IMP,
)
print("python-pptx:", r2.stdout.strip() if r2.returncode == 0 else f"NG: {r2.stderr.strip()[:80]}")

# 3. 全形式のparse_file分岐確認
r3 = subprocess.run(
    [
        sys.executable,
        "-c",
        """
import sys; sys.path.insert(0, '.')
import inspect
from file_parser import parse_file, parse_pptx, parse_csv
src = inspect.getsource(parse_file)
checks = {
    "xlsx/xls": ".xlsx" in src and ".xls" in src,
    "pdf":      ".pdf" in src,
    "docx/doc": ".docx" in src and ".doc" in src,
    "pptx/ppt": ".pptx" in src and ".ppt" in src,
    "csv/tsv":  ".csv" in src and ".tsv" in src,
    "parse_pptx存在": True,
    "parse_csv存在":  True,
}
for k, v in checks.items():
    print(f"  {k}: {'OK' if v else 'NG'}")
""",
    ],
    capture_output=True,
    encoding="utf-8",
    errors="replace",
    cwd=IMP,
)
print("分岐チェック:")
print(r3.stdout if r3.returncode == 0 else f"NG: {r3.stderr.strip()[-200:]}")

# 4. pptxダミーファイルで動作テスト
r4 = subprocess.run(
    [
        sys.executable,
        "-c",
        """
import sys; sys.path.insert(0, '.')
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# ダミーpptx生成
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
tf = txBox.text_frame
tf.text = "氏名: テスト太郎\\n単価: 70万円\\nスキル: Python, AWS"
buf = BytesIO()
prs.save(buf)
data = buf.getvalue()

from file_parser import parse_file
result = parse_file("test.pptx", ".pptx", data)
print("pptx parse OK:", result[:80] if result else "None")
""",
    ],
    capture_output=True,
    encoding="utf-8",
    errors="replace",
    cwd=IMP,
)
print("pptx動作テスト:", r4.stdout.strip() if r4.returncode == 0 else f"NG: {r4.stderr.strip()[-200:]}")

# 5. csvダミーで動作テスト
r5 = subprocess.run(
    [
        sys.executable,
        "-c",
        """
import sys; sys.path.insert(0, '.')
from file_parser import parse_file
csv_data = "氏名,単価,スキル\\nテスト太郎,70,Python\\n".encode("utf-8")
result = parse_file("test.csv", ".csv", csv_data)
print("csv parse OK:", result[:80] if result else "None")
""",
    ],
    capture_output=True,
    encoding="utf-8",
    errors="replace",
    cwd=IMP,
)
print("csv動作テスト:", r5.stdout.strip() if r5.returncode == 0 else f"NG: {r5.stderr.strip()[-200:]}")
