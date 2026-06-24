import sys

sys.stdout.reconfigure(encoding="utf-8")

fp_path = r"C:\Users\ma_py\OneDrive\デスクトップ\ses_work\mail_attachment_importer\file_parser.py"

with open(fp_path, encoding="utf-8") as f:
    content = f.read()

# バックアップ
with open(fp_path + ".bak_pptx_csv", "w", encoding="utf-8") as f:
    f.write(content)
print("backup: OK")

# --- 追加する関数 ---
new_functions = '''

def parse_pptx(data: bytes) -> str:
    """python-pptxで全スライドのテキストを抽出（テーブル含む）"""
    try:
        from pptx import Presentation
        from io import BytesIO
        prs = Presentation(BytesIO(data))
        lines = []
        for i, slide in enumerate(prs.slides, start=1):
            lines.append(f"=== スライド {i} ===")
            for shape in slide.shapes:
                # テキストフレーム
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
                # テーブル
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            lines.append("\\t".join(cells))
        return "\\n".join(lines)
    except Exception as e:
        logger.error(f"PowerPoint解析失敗: {e}")
        raise


def parse_csv(data: bytes) -> str:
    """CSVをタブ区切りテキストに変換（UTF-8→cp932フォールバック）"""
    import csv
    from io import StringIO
    for encoding in ('utf-8-sig', 'utf-8', 'cp932', 'shift_jis'):
        try:
            text = data.decode(encoding)
            reader = csv.reader(StringIO(text))
            lines = []
            for row in reader:
                cells = [cell.strip() for cell in row if cell.strip()]
                if cells:
                    lines.append("\\t".join(cells))
            return "\\n".join(lines)
        except (UnicodeDecodeError, Exception):
            continue
    logger.error("CSV解析失敗: エンコード特定不能")
    raise ValueError("CSV encoding not detected")

'''

# parse_file()の前に挿入
insert_before = "\ndef parse_file("
if "parse_pptx" not in content:
    content = content.replace(insert_before, new_functions + insert_before, 1)
    print("関数追加: OK")
else:
    print("関数: 既に存在")

# parse_file()の分岐拡張
old_else = """        else:
            logger.warning(f"対応外形式: {ext} ({filename})")
            return None"""

new_else = """        elif ext in (".pptx", ".ppt"):
            return parse_pptx(data)
        elif ext in (".csv", ".tsv"):
            return parse_csv(data)
        else:
            logger.warning(f"対応外形式: {ext} ({filename})")
            return None"""

if ".pptx" not in content:
    content = content.replace(old_else, new_else, 1)
    print("parse_file分岐追加: OK")
else:
    print("parse_file分岐: 既に存在")

with open(fp_path, "w", encoding="utf-8") as f:
    f.write(content)
print("ファイル書き込み: OK")
